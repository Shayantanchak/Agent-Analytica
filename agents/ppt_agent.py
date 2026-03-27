from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from io import BytesIO
import tempfile
import os

class PPTAgent:
    """
    Phase 4: Premium PPT Output Generation
    Creates highly styled, boardroom-ready automated PPTX slides.
    Colors based on Analytica theme: Dark (#0F172A) background, Blue accents (#2563EB).
    """
    def __init__(self):
        # Brand Colors
        self.bg_color = RGBColor(15, 23, 42)      # #0F172A Dark background
        self.card_color = RGBColor(30, 41, 59)    # #1E293B Cards
        self.accent_color = RGBColor(37, 99, 235) # #2563EB Blue accent
        self.text_light = RGBColor(255, 255, 255) # White
        self.text_gray = RGBColor(173, 181, 189)  # Light gray

    def _set_background(self, slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color
        
    def _add_header(self, slide, title_text):
        # Accent gradient simulation (top border line)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = self.accent_color
        line.line.color.rgb = self.accent_color
        
        # Title text
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.text_light

    def execute(self, analysis_results: dict) -> BytesIO:
        # Use a blank presentation
        prs = Presentation()
        # Set to Widescreen 16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        blank_layout = prs.slide_layouts[6] # completely blank template
        
        # --- Slide 1: Welcome / Title Slide ---
        slide1 = prs.slides.add_slide(blank_layout)
        self._set_background(slide1)
        
        # Big Center Title
        txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Analytica Intelligence"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.text_light
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        txBox_sub = slide1.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
        tf_sub = txBox_sub.text_frame
        p_sub = tf_sub.add_paragraph()
        p_sub.text = "Automated Executive Data Insights & KPIs"
        p_sub.font.size = Pt(20)
        p_sub.font.color.rgb = self.accent_color
        p_sub.alignment = PP_ALIGN.CENTER
        
        # --- Slide 2: Executive Summary & KPIs ---
        slide2 = prs.slides.add_slide(blank_layout)
        self._set_background(slide2)
        self._add_header(slide2, "Executive KPI Summary")
        
        # Draw KPIs as Beautiful "Cards" on the slide
        kpis = analysis_results.get("kpis", {})
        if kpis:
            start_x = 0.5
            y_pos = 1.3
            box_width = 2.8
            box_height = 1.2
            spacing = 0.2
            
            for i, (k, v) in enumerate(kpis.items()):
                # Wrap to next line if more than 3 per row
                x_pos = start_x + (i % 3) * (box_width + spacing)
                if i > 0 and i % 3 == 0:
                    y_pos += box_height + spacing
                
                # Draw Card Box
                shape = slide2.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, 
                    Inches(x_pos), Inches(y_pos), Inches(box_width), Inches(box_height)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.card_color
                shape.line.color.rgb = self.accent_color
                
                # KPI Title
                tx = slide2.shapes.add_textbox(Inches(x_pos), Inches(y_pos+0.1), Inches(box_width), Inches(0.4))
                p = tx.text_frame.paragraphs[0]
                p.text = str(k).upper()
                p.font.size = Pt(12)
                p.font.color.rgb = self.text_gray
                p.alignment = PP_ALIGN.CENTER
                
                # KPI Value
                tx2 = slide2.shapes.add_textbox(Inches(x_pos), Inches(y_pos+0.4), Inches(box_width), Inches(0.6))
                p2 = tx2.text_frame.paragraphs[0]
                p2.text = str(v)
                p2.font.size = Pt(28)
                p2.font.bold = True
                p2.font.color.rgb = self.text_light
                p2.alignment = PP_ALIGN.CENTER

        # AI Narrative below the KPI Cards
        narrative = analysis_results.get("narrative", "").replace("### ", "").replace("**", "")
        if narrative:
            txBox_n = slide2.shapes.add_textbox(Inches(0.5), Inches(y_pos + box_height + 0.3), Inches(9), Inches(2))
            lines = [line for line in narrative.split('\n') if line.strip() and not line.startswith('>')]
            tf_n = txBox_n.text_frame
            tf_n.word_wrap = True
            for line in lines[:4]: # Grab up to 4 bullet points safely
                p = tf_n.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                p.font.color.rgb = self.text_light
                p.level = 0
            
        # --- Slide 3+: Dynamic AI Charts ---
        if "charts" in analysis_results:
            for chart_name, fig in analysis_results["charts"].items():
                try:
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                        img_path = tmp.name
                        
                    # Inject a matching Dark-Mode theme directly into the Plotly chart export!
                    f_copy = fig
                    f_copy.update_layout(
                        paper_bgcolor="#0F172A", # Master PPT Background color
                        plot_bgcolor="#0F172A", 
                        font={'color': 'white'}
                    )
                    f_copy.write_image(img_path, engine="kaleido", width=900, height=400)
                    
                    slide_c = prs.slides.add_slide(blank_layout)
                    self._set_background(slide_c)
                    self._add_header(slide_c, f"Visual Query: {chart_name.replace('_', ' ').title()}")
                    
                    # Insert the dynamically themed Image perfectly centered
                    slide_c.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(9))
                    os.unlink(img_path)
                except Exception as e:
                    print(f"Failed to write image for {chart_name}: {e}")
                    
        # Return object as byte stream for immediate UI download
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        return output
